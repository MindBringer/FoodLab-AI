import json
import os
import subprocess
import tempfile
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Any, Dict, List, Optional

import extract_msg
import mailparser
from bs4 import BeautifulSoup
from docx import Document
from fastapi import FastAPI, File, HTTPException, UploadFile
from openpyxl import load_workbook
from pptx import Presentation
from pydantic import BaseModel
from pypdf import PdfReader

app = FastAPI(title="FoodLab Parser Service", version="3.0.0")

DATA_DIR = Path(os.getenv("DATA_DIR", "/srv/foodlab/data"))
OCR_DIR = DATA_DIR / "ocr"
for p in [DATA_DIR, OCR_DIR]:
    p.mkdir(parents=True, exist_ok=True)

class ParseRequest(BaseModel):
    file_path: str
    mime_type: Optional[str] = None
    force_ocr: bool = False

def _read_pdf(path: Path) -> Dict[str, Any]:
    reader = PdfReader(str(path))
    texts = []
    for page_no, page in enumerate(reader.pages, start=1):
        texts.append({"page": page_no, "text": page.extract_text() or ""})
    text = "\n\n".join([x["text"] for x in texts]).strip()
    return {"text": text, "pages": texts, "metadata": {"page_count": len(texts)}}

def _ocr_pdf(path: Path) -> Path:
    out_path = OCR_DIR / f"{path.stem}.ocr.pdf"
    cmd = ["ocrmypdf", "--force-ocr", "--skip-text", str(path), str(out_path)]
    subprocess.run(cmd, check=True, capture_output=True)
    return out_path

def _read_docx(path: Path) -> Dict[str, Any]:
    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return {"text": "\n".join(paragraphs), "sections": paragraphs}

def _read_xlsx(path: Path) -> Dict[str, Any]:
    wb = load_workbook(str(path), data_only=True)
    sheets = []
    parts = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            values = [str(v) for v in row if v is not None and str(v).strip()]
            if values:
                rows.append(values)
                parts.append(" | ".join(values))
        sheets.append({"name": ws.title, "rows": rows})
    return {"text": "\n".join(parts), "sheets": sheets}

def _read_pptx(path: Path) -> Dict[str, Any]:
    prs = Presentation(str(path))
    slides = []
    parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                val = shape.text.strip()
                if val:
                    texts.append(val)
                    parts.append(val)
        slides.append({"slide": idx, "text": "\n".join(texts)})
    return {"text": "\n\n".join(parts), "slides": slides}

def _read_txt(path: Path) -> Dict[str, Any]:
    return {"text": path.read_text(encoding="utf-8", errors="ignore")}

def _read_html(path: Path) -> Dict[str, Any]:
    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "lxml")
    return {"text": soup.get_text("\n", strip=True)}

def _read_eml(path: Path) -> Dict[str, Any]:
    parsed = mailparser.parse_from_file(str(path))
    parts = [
        parsed.subject or "",
        parsed.body or "",
        "\n".join(parsed.text_plain or []),
        "\n".join(parsed.text_html or []),
    ]
    attachments = [{"filename": a.get("filename"), "mail_content_type": a.get("mail_content_type")} for a in parsed.attachments]
    return {
        "text": "\n\n".join([p for p in parts if p.strip()]),
        "email": {
            "subject": parsed.subject,
            "from": parsed.from_,
            "to": parsed.to,
            "date": parsed.date.isoformat() if parsed.date else None,
            "attachments": attachments,
        },
    }

def _read_msg(path: Path) -> Dict[str, Any]:
    msg = extract_msg.Message(str(path))
    parts = [msg.subject or "", msg.body or ""]
    return {
        "text": "\n\n".join([p for p in parts if p.strip()]),
        "email": {
            "subject": msg.subject,
            "from": msg.sender,
            "to": msg.to,
            "date": msg.date,
        },
    }

def parse_path(path: Path, force_ocr: bool = False) -> Dict[str, Any]:
    suffix = path.suffix.lower()
    result: Dict[str, Any]
    ocr_applied = False

    if suffix == ".pdf":
        result = _read_pdf(path)
        if force_ocr or not (result.get("text") or "").strip():
            ocr_path = _ocr_pdf(path)
            result = _read_pdf(ocr_path)
            ocr_applied = True
    elif suffix == ".docx":
        result = _read_docx(path)
    elif suffix == ".xlsx":
        result = _read_xlsx(path)
    elif suffix == ".pptx":
        result = _read_pptx(path)
    elif suffix in {".txt", ".md", ".csv", ".json", ".log"}:
        result = _read_txt(path)
    elif suffix in {".htm", ".html"}:
        result = _read_html(path)
    elif suffix == ".eml":
        result = _read_eml(path)
    elif suffix == ".msg":
        result = _read_msg(path)
    else:
        raise HTTPException(status_code=415, detail=f"unsupported file type: {suffix}")

    return {
        "file_path": str(path),
        "source_name": path.name,
        "source_type": suffix.lstrip("."),
        "ocr_applied": ocr_applied,
        **result,
    }

@app.get("/health")
def health():
    return {"status": "ok", "service": "parser-service"}

@app.post("/parse")
def parse(req: ParseRequest):
    path = Path(req.file_path)
    if not path.exists():
        raise HTTPException(status_code=404, detail="file not found")
    return parse_path(path, force_ocr=req.force_ocr)

@app.post("/parse/upload")
async def parse_upload(file: UploadFile = File(...), force_ocr: bool = False):
    suffix = Path(file.filename or "upload.bin").suffix
    temp_path = DATA_DIR / "inbox" / f"upload_{os.getpid()}_{Path(file.filename or 'upload').stem}{suffix}"
    temp_path.parent.mkdir(parents=True, exist_ok=True)
    temp_path.write_bytes(await file.read())
    return parse_path(temp_path, force_ocr=force_ocr)
